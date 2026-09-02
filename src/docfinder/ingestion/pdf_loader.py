"""Document loading and chunking utilities.

Supports PDF (via PyMuPDF), plain text (.txt), Markdown (.md),
Word documents (.docx via python-docx, .doc via olefile),
OpenDocument formats (.odt/.odp/.odg via odfpy),
PowerPoint (.pptx via python-pptx, .ppt via olefile),
HTML (.html/.htm via beautifulsoup4),
and EPUB (.epub via zipfile + beautifulsoup4).
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Dict, Iterable, Iterator

import fitz  # PyMuPDF

from docfinder.models import ChunkRecord
from docfinder.utils.text import normalize_whitespace

LOGGER = logging.getLogger(__name__)


# ── PDF ───────────────────────────────────────────────────────────────────────


def _table_to_markdown(table) -> str:
    """Convert a PyMuPDF Table object to a Markdown table string."""
    rows = table.extract()
    if not rows:
        return ""

    lines: list[str] = []
    # Header row
    header = [str(cell or "").strip().replace("\n", " ") for cell in rows[0]]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    # Data rows
    for row in rows[1:]:
        cells = [str(cell or "").strip().replace("\n", " ") for cell in row]
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def _extract_page_text(page) -> str:
    """Extract text from a PDF page, converting tables to Markdown.

    If tables are detected, they are rendered as Markdown and inserted
    in place of the raw table text.  Non-table text is extracted normally.
    """
    try:
        tables = page.find_tables()
    except Exception:
        tables = None

    if not tables or not tables.tables:
        # No tables — fast path, same as before
        return page.get_text() or ""

    # Collect table bounding boxes and their Markdown representations
    table_items: list[tuple[float, str]] = []  # (top_y, markdown)
    table_rects: list[fitz.Rect] = []
    for tab in tables.tables:
        md = _table_to_markdown(tab)
        if md:
            rect = fitz.Rect(tab.bbox)
            table_items.append((rect.y0, md))
            table_rects.append(rect)

    # Extract text blocks excluding table areas
    text_blocks: list[tuple[float, str]] = []  # (top_y, text)
    for block in page.get_text("blocks") or []:
        # block = (x0, y0, x1, y1, text, block_no, block_type)
        if block[6] != 0:  # skip image blocks
            continue
        block_rect = fitz.Rect(block[:4])
        # Skip blocks that overlap significantly with any table
        overlaps_table = False
        for tr in table_rects:
            intersection = block_rect & tr
            if not intersection.is_empty:
                block_area = block_rect.width * block_rect.height
                if block_area > 0:
                    overlap_ratio = (intersection.width * intersection.height) / block_area
                    if overlap_ratio > 0.5:
                        overlaps_table = True
                        break
        if not overlaps_table:
            text = block[4].strip()
            if text:
                text_blocks.append((block[1], text))

    # Merge text blocks and tables, sorted by vertical position
    all_parts: list[tuple[float, str]] = text_blocks + table_items
    all_parts.sort(key=lambda x: x[0])

    return "\n\n".join(part[1] for part in all_parts)


def iter_text_parts(path: Path) -> Iterator[str]:
    """Yield text content from a PDF file, page by page."""
    try:
        doc = fitz.open(path)
    except Exception as exc:
        LOGGER.error("Failed to open PDF %s: %s", path, exc)
        return

    try:
        for index in range(len(doc)):
            try:
                page = doc[index]
                text = _extract_page_text(page)
                normalized = normalize_whitespace([text])
                if normalized:
                    yield normalized + "\n"
            except Exception as exc:  # pragma: no cover
                LOGGER.warning("Failed to read page %s in %s: %s", index, path, exc)
    finally:
        doc.close()


def get_pdf_metadata(path: Path) -> Dict[str, str]:
    """Extract title and page count from a PDF."""
    doc = fitz.open(path)
    try:
        metadata = doc.metadata or {}
        title = metadata.get("title") or path.stem
        return {"title": title, "page_count": str(len(doc))}
    finally:
        doc.close()


# ── Plain text ────────────────────────────────────────────────────────────────

# Virtual page size for plain text files (characters per virtual page).
_TXT_VIRTUAL_PAGE_CHARS = 3000


def iter_text_parts_txt(path: Path) -> Iterator[str]:
    """Yield the full content of a .txt file as a single string."""
    try:
        yield path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_txt_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(virtual_page, text)`` for plain text files.

    Splits the file into virtual pages of ~3000 characters each,
    breaking at the nearest newline to avoid cutting mid-sentence.
    """
    try:
        content = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)
        return

    if not content.strip():
        return

    page = 1
    start = 0
    while start < len(content):
        end = start + _TXT_VIRTUAL_PAGE_CHARS
        if end < len(content):
            # Try to break at a newline within the last 20% of the slice
            search_start = max(start, end - _TXT_VIRTUAL_PAGE_CHARS // 5)
            nl = content.rfind("\n", search_start, end)
            if nl > start:
                end = nl + 1
        chunk = content[start:end]
        if chunk.strip():
            yield page, chunk
        page += 1
        start = end


# ── Markdown ──────────────────────────────────────────────────────────────────

_MD_HEADING_LINE = re.compile(r"^#{1,6}\s+", re.MULTILINE)
_MD_HEADING_SPLIT = re.compile(r"(?=^#{1,6}\s)", re.MULTILINE)
_MD_BOLD_ITAL = re.compile(r"\*{1,3}([^*\n]+)\*{1,3}")
_MD_CODE = re.compile(r"`{1,3}[^`]*`{1,3}", re.DOTALL)
_MD_IMAGE = re.compile(r"!\[.*?\]\([^)]*\)")
_MD_LINK = re.compile(r"\[([^\]]+)\]\([^)]+\)")
_MD_HR = re.compile(r"^[-*_]{3,}\s*$", re.MULTILINE)


def _clean_md(text: str) -> str:
    """Strip Markdown formatting from a text block."""
    text = _MD_CODE.sub("", text)
    text = _MD_IMAGE.sub("", text)
    text = _MD_LINK.sub(r"\1", text)
    text = _MD_BOLD_ITAL.sub(r"\1", text)
    text = _MD_HR.sub("", text)
    text = _MD_HEADING_LINE.sub("", text)
    return text


def iter_text_parts_md(path: Path) -> Iterator[str]:
    """Yield plain text extracted from a Markdown file."""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        yield _clean_md(text)
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_md_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(section_number, text)`` for Markdown files.

    Each top-level heading (``# …``, ``## …``, etc.) starts a new
    virtual page.  Content before the first heading is section 1.
    """
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)
        return

    sections = _MD_HEADING_SPLIT.split(raw)
    page = 1
    for section in sections:
        cleaned = _clean_md(section).strip()
        if cleaned:
            yield page, cleaned + "\n"
            page += 1


# ── Word (.docx) ──────────────────────────────────────────────────────────────

# Number of paragraphs per virtual page for .docx files.
_DOCX_PARAS_PER_PAGE = 10


def _import_docx_document():
    """Import and return the ``Document`` class from python-docx, or *None*."""
    try:
        from docx import Document  # type: ignore[import-untyped]

        return Document
    except ImportError:
        LOGGER.warning(
            "python-docx not installed — cannot index .docx files. "
            "Install with: pip install python-docx"
        )
        return None


def iter_text_parts_docx(path: Path) -> Iterator[str]:
    """Yield paragraph text from a .docx file."""
    Document = _import_docx_document()
    if Document is None:
        return
    try:
        doc = Document(str(path))
        for para in doc.paragraphs:
            stripped = para.text.strip()
            if stripped:
                yield stripped + "\n"
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_docx_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(virtual_page, text)`` for Word documents.

    Groups every 10 consecutive non-empty paragraphs into one virtual page.
    """
    Document = _import_docx_document()
    if Document is None:
        return
    try:
        doc = Document(str(path))
        buf: list[str] = []
        page = 1
        for para in doc.paragraphs:
            stripped = para.text.strip()
            if not stripped:
                continue
            buf.append(stripped)
            if len(buf) >= _DOCX_PARAS_PER_PAGE:
                yield page, "\n".join(buf) + "\n"
                buf = []
                page += 1
        if buf:
            yield page, "\n".join(buf) + "\n"
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


# ── ODF (ODT / ODP / ODG) ─────────────────────────────────────────────────────

_ODF_PARAS_PER_PAGE = 10


def _import_odf():
    """Import odfpy, return (load, P_element_class, extractText) or (None, None, None)."""
    try:
        from odf import teletype
        from odf.opendocument import load
        from odf.text import P

        return load, P, teletype.extractText
    except ImportError:
        LOGGER.warning(
            "odfpy not installed — cannot index ODF files (.odt/.odp/.odg). "
            "Install with: pip install odfpy"
        )
        return None, None, None


def iter_text_parts_odf(path: Path) -> Iterator[str]:
    """Yield text content from an OpenDocument file."""
    load, P, extract_text = _import_odf()
    if load is None:
        return
    try:
        doc = load(str(path))
        for p_elem in doc.getElementsByType(P):
            text = extract_text(p_elem)
            if text.strip():
                yield text.strip() + "\n"
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_odf_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(virtual_page, text)`` grouping ~10 paragraphs."""
    load, P, extract_text = _import_odf()
    if load is None:
        return
    try:
        doc = load(str(path))
        buf: list[str] = []
        page = 1
        for p_elem in doc.getElementsByType(P):
            text = extract_text(p_elem).strip()
            if not text:
                continue
            buf.append(text)
            if len(buf) >= _ODF_PARAS_PER_PAGE:
                yield page, "\n".join(buf) + "\n"
                buf = []
                page += 1
        if buf:
            yield page, "\n".join(buf) + "\n"
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


# ── PPTX ──────────────────────────────────────────────────────────────────────


def _import_pptx():
    """Import and return the ``Presentation`` class from python-pptx, or None."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        return Presentation
    except ImportError:
        LOGGER.warning(
            "python-pptx not installed — cannot index .pptx files. "
            "Install with: pip install python-pptx"
        )
        return None


def _extract_pptx_slide_text(slide) -> str:
    """Extract all text from a single pptx Slide object."""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
    return "\n".join(parts)


def iter_text_parts_pptx(path: Path) -> Iterator[str]:
    """Yield text from every slide in a .pptx file."""
    Presentation = _import_pptx()
    if Presentation is None:
        return
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            text = _extract_pptx_slide_text(slide)
            if text:
                yield text + "\n"
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_pptx_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(slide_number, text)`` — one page per slide."""
    Presentation = _import_pptx()
    if Presentation is None:
        return
    try:
        prs = Presentation(str(path))
        for page_num, slide in enumerate(prs.slides, 1):
            text = _extract_pptx_slide_text(slide)
            if text:
                yield page_num, text + "\n"
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


# ── HTML ──────────────────────────────────────────────────────────────────────

_HTML_VIRTUAL_PAGE_CHARS = 3000


def _import_beautifulsoup4():
    """Import and return ``BeautifulSoup``, or None."""
    try:
        from bs4 import BeautifulSoup  # type: ignore[import-untyped]

        return BeautifulSoup
    except ImportError:
        LOGGER.warning(
            "beautifulsoup4 not installed — cannot index .html/.epub files. "
            "Install with: pip install beautifulsoup4"
        )
        return None


def iter_text_parts_html(path: Path) -> Iterator[str]:
    """Yield cleaned plain text from an HTML file."""
    BeautifulSoup = _import_beautifulsoup4()
    if BeautifulSoup is None:
        return
    try:
        raw = path.read_bytes()
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        if text:
            yield text
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_html_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(virtual_page, text)`` for HTML files (~3000 chars per page)."""
    BeautifulSoup = _import_beautifulsoup4()
    if BeautifulSoup is None:
        return
    try:
        raw = path.read_bytes()
        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        if not text:
            return

        page = 1
        start = 0
        while start < len(text):
            end = start + _HTML_VIRTUAL_PAGE_CHARS
            if end < len(text):
                search_start = max(start, end - _HTML_VIRTUAL_PAGE_CHARS // 5)
                nl = text.rfind("\n", search_start, end)
                if nl > start:
                    end = nl + 1
            chunk = text[start:end]
            if chunk.strip():
                yield page, chunk + "\n"
            page += 1
            start = end
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


# ── EPUB ──────────────────────────────────────────────────────────────────────


def _parse_epub(path: Path, BeautifulSoup) -> Iterator[tuple[int, str]]:
    """Parse an EPUB and yield ``(chapter_number, text)`` from spine items."""
    import zipfile
    from xml.etree import ElementTree

    with zipfile.ZipFile(path) as zf:
        try:
            container_data = zf.read("META-INF/container.xml")
        except KeyError:
            return

        ns_ct = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
        container = ElementTree.fromstring(container_data)
        rootfile = container.find(".//c:rootfile", ns_ct)
        if rootfile is None:
            return
        opf_rel = rootfile.get("full-path", "")
        if not opf_rel:
            return

        base_dir = opf_rel.rsplit("/", 1)[0] if "/" in opf_rel else ""

        opf_data = zf.read(opf_rel)
        opf = ElementTree.fromstring(opf_data)

        ns_opf = {"opf": "http://www.idpf.org/2007/opf"}
        manifest: dict[str, str] = {}
        for item in opf.findall(".//opf:manifest/opf:item", ns_opf):
            item_id = item.get("id", "")
            href = item.get("href", "")
            if href and base_dir:
                href = base_dir + "/" + href
            manifest[item_id] = href

        spine = opf.find(".//opf:spine", ns_opf)
        if spine is None:
            return

        page_num = 0
        for ref in spine.findall("opf:itemref", ns_opf):
            idref = ref.get("idref", "")
            href = manifest.get(idref, "")
            if not href:
                continue
            try:
                content = zf.read(href)
            except KeyError:
                continue

            soup = BeautifulSoup(content, "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
            if text:
                page_num += 1
                yield page_num, text + "\n"


def iter_text_parts_epub(path: Path) -> Iterator[str]:
    """Yield chapter text from an EPUB file."""
    BeautifulSoup = _import_beautifulsoup4()
    if BeautifulSoup is None:
        return
    try:
        for _, text in _parse_epub(path, BeautifulSoup):
            yield text
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


def iter_text_parts_epub_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(chapter_number, text)`` — one page per spine item."""
    BeautifulSoup = _import_beautifulsoup4()
    if BeautifulSoup is None:
        return
    try:
        yield from _parse_epub(path, BeautifulSoup)
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)


# ── DOC (Word 97‑2003 binary) ────────────────────────────────────────────────

_DOC_VIRTUAL_PAGE_CHARS = 3000


def _import_olefile():
    """Import and return the ``olefile`` module, or None."""
    try:
        import olefile  # type: ignore[import-untyped]

        return olefile
    except ImportError:
        LOGGER.warning(
            "olefile not installed — cannot index .doc/.ppt files. "
            "Install with: pip install olefile"
        )
        return None


def _extract_text_from_doc_binary(path: Path) -> Iterator[str]:
    """Extract text from a Word 97-2003 file using its piece table."""
    olefile = _import_olefile()
    if olefile is None:
        return
    try:
        ole = olefile.OleFileIO(path)
        try:
            word_data = ole.openstream("WordDocument").read()
            if len(word_data) < 32 or word_data[:2] not in (
                b"\xec\xa5",  # Word 97-2003
                b"\xdc\xa5",  # Word 6/95
                b"\xdb\xa5",  # Word 95
            ):
                LOGGER.warning("Unsupported or invalid Word binary file: %s", path)
                return

            csw = int.from_bytes(word_data[32:34], "little")
            offset = 34 + csw * 2
            cslw = int.from_bytes(word_data[offset : offset + 2], "little")
            offset += 2 + cslw * 4
            cb_rg_fc_lcb = int.from_bytes(word_data[offset : offset + 2], "little")
            offset += 2
            fc_lcb = word_data[offset : offset + cb_rg_fc_lcb * 8]
            # fcClx is entry 33 in FibRgFcLcb97 (fc followed by lcb).
            clx_entry = 33 * 8
            if len(fc_lcb) < clx_entry + 8:
                LOGGER.warning("Word binary file has no CLX reference: %s", path)
                return
            fc_clx = int.from_bytes(fc_lcb[clx_entry : clx_entry + 4], "little")
            lcb_clx = int.from_bytes(fc_lcb[clx_entry + 4 : clx_entry + 8], "little")
            if not lcb_clx:
                return

            flags = int.from_bytes(word_data[10:12], "little")
            table_name = "1Table" if flags & 0x0200 else "0Table"
            table_data = ole.openstream(table_name).read()
            clx = table_data[fc_clx : fc_clx + lcb_clx]
        finally:
            ole.close()
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)
        return

    # CLX may begin with formatting records before the piece table.
    pos = 0
    while pos < len(clx) and clx[pos] == 0x01:
        if pos + 5 > len(clx):
            return
        pos += 5 + int.from_bytes(clx[pos + 1 : pos + 5], "little")
    if pos >= len(clx) or clx[pos] != 0x02 or pos + 5 > len(clx):
        LOGGER.warning("Word binary file has an invalid piece table: %s", path)
        return

    piece_table_size = int.from_bytes(clx[pos + 1 : pos + 5], "little")
    piece_table = clx[pos + 5 : pos + 5 + piece_table_size]
    if len(piece_table) < 4 or (len(piece_table) - 4) % 12:
        LOGGER.warning("Word binary file has a truncated piece table: %s", path)
        return

    piece_count = (len(piece_table) - 4) // 12
    cps = [
        int.from_bytes(piece_table[index * 4 : index * 4 + 4], "little")
        for index in range(piece_count + 1)
    ]
    text_parts: list[str] = []
    for index in range(piece_count):
        pcd_offset = 4 * (piece_count + 1) + index * 8
        fc = int.from_bytes(piece_table[pcd_offset + 2 : pcd_offset + 6], "little")
        cp_length = cps[index + 1] - cps[index]
        if cp_length <= 0:
            continue
        compressed = bool(fc & 0x40000000)
        if compressed:
            # Compressed pieces store the byte offset multiplied by two.
            start = (fc & 0x3FFFFFFF) // 2
            raw = word_data[start : start + cp_length]
            text_parts.append(raw.decode("cp1252", errors="replace"))
        else:
            start = fc
            raw = word_data[start : start + cp_length * 2]
            text_parts.append(raw.decode("utf-16-le", errors="replace"))

    text = "".join(text_parts)
    text = text.replace("\x07", "\n").replace("\x0b", "\n").replace("\x0c", "\n")
    text = "".join(char for char in text if char in "\n\r\t" or char.isprintable())
    text = re.sub(r"[^\S\n\r\t]+", " ", text).strip()
    if text:
        yield text


def iter_text_parts_doc(path: Path) -> Iterator[str]:
    """Yield text extracted from a .doc file."""
    yield from _extract_text_from_doc_binary(path)


def iter_text_parts_doc_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(virtual_page, text)`` for Word binary documents."""
    text = "\n".join(_extract_text_from_doc_binary(path))
    if not text.strip():
        return

    page = 1
    start = 0
    while start < len(text):
        end = start + _DOC_VIRTUAL_PAGE_CHARS
        if end < len(text):
            search_start = max(start, end - _DOC_VIRTUAL_PAGE_CHARS // 5)
            nl = text.rfind("\n", search_start, end)
            if nl > start:
                end = nl + 1
        chunk = text[start:end]
        if chunk.strip():
            yield page, chunk + "\n"
        page += 1
        start = end


# ── PPT (PowerPoint 97‑2003 binary) ───────────────────────────────────────────

_PPT_VIRTUAL_PAGE_CHARS = 3000

# Record types for the PowerPoint binary format
_PPT_TEXT_CHARS_ATOM = 0x0FA0
_PPT_TEXT_BYTES_ATOM = 0x0FA8


def _walk_ppt_records(data: bytes, offset: int, end: int, texts: list[str]) -> None:
    """Walk a PowerPoint binary record tree looking for text atoms.

    Record headers are little-endian ([MS-PPT]): ver+instance (u16), type
    (u16), length (u32). A record is a container when its version nibble
    is ``0xF``.
    """
    import struct

    while offset + 8 <= end:
        # rec_ver = data[offset] & 0x0F  (used below for container detection)
        # rec_instance = ((data[offset] >> 4) & 0x0F) << 8 | data[offset + 1]  (not needed)
        rec_type = struct.unpack_from("<H", data, offset + 2)[0]
        rec_len = struct.unpack_from("<I", data, offset + 4)[0]

        data_start = offset + 8
        data_end = data_start + rec_len

        if data_end > end:
            break

        if rec_type == _PPT_TEXT_CHARS_ATOM:
            try:
                raw = data[data_start:data_end].decode("utf-16-le", errors="replace")
                clean = "".join(c if c.isprintable() or c in "\n\r\t" else " " for c in raw)
                clean = re.sub(r"\s+", " ", clean).strip()
                if clean:
                    texts.append(clean)
            except Exception:
                pass
        elif rec_type == _PPT_TEXT_BYTES_ATOM:
            try:
                raw = data[data_start:data_end].decode("utf-8", errors="replace")
                clean = "".join(c if c.isprintable() or c in "\n\r\t" else " " for c in raw)
                clean = re.sub(r"\s+", " ", clean).strip()
                if clean:
                    texts.append(clean)
            except Exception:
                pass

        if (data[offset] & 0x0F) == 0x0F:
            _walk_ppt_records(data, data_start, data_end, texts)

        offset = data_end


def _extract_text_from_ppt_binary(path: Path) -> Iterator[str]:
    """Extract text from PowerPoint 97‑2003 (.ppt) files via olefile."""
    olefile = _import_olefile()
    if olefile is None:
        return
    try:
        ole = olefile.OleFileIO(path)
        try:
            stream = ole.openstream("PowerPoint Document")
            data = stream.read()
        finally:
            ole.close()
    except Exception as exc:
        LOGGER.error("Failed to read %s: %s", path, exc)
        return

    texts: list[str] = []
    _walk_ppt_records(data, 0, len(data), texts)

    for text in texts:
        if text.strip():
            yield text.strip() + "\n"


def iter_text_parts_ppt(path: Path) -> Iterator[str]:
    """Yield text from a .ppt file."""
    yield from _extract_text_from_ppt_binary(path)


def iter_text_parts_ppt_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(virtual_page, text)`` for PowerPoint binary documents."""
    text = "\n".join(_extract_text_from_ppt_binary(path))
    if not text.strip():
        return

    page = 1
    start = 0
    while start < len(text):
        end = start + _PPT_VIRTUAL_PAGE_CHARS
        if end < len(text):
            search_start = max(start, end - _PPT_VIRTUAL_PAGE_CHARS // 5)
            nl = text.rfind("\n", search_start, end)
            if nl > start:
                end = nl + 1
        chunk = text[start:end]
        if chunk.strip():
            yield page, chunk + "\n"
        page += 1
        start = end


# ── Dispatcher ────────────────────────────────────────────────────────────────


def _iter_text_by_suffix(path: Path) -> Iterator[str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        yield from iter_text_parts(path)
    elif suffix == ".txt":
        yield from iter_text_parts_txt(path)
    elif suffix == ".md":
        yield from iter_text_parts_md(path)
    elif suffix == ".docx":
        yield from iter_text_parts_docx(path)
    elif suffix in (".odt", ".odp", ".odg"):
        yield from iter_text_parts_odf(path)
    elif suffix == ".pptx":
        yield from iter_text_parts_pptx(path)
    elif suffix in (".html", ".htm"):
        yield from iter_text_parts_html(path)
    elif suffix == ".epub":
        yield from iter_text_parts_epub(path)
    elif suffix == ".doc":
        yield from iter_text_parts_doc(path)
    elif suffix == ".ppt":
        yield from iter_text_parts_ppt(path)
    else:
        LOGGER.warning("Unsupported file type: %s", path.suffix)


def _get_title(path: Path) -> str:
    if path.suffix.lower() == ".pdf":
        try:
            return get_pdf_metadata(path).get("title", path.stem)
        except Exception:
            pass
    return path.stem


def iter_text_parts_paged(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(page_number, text)`` for PDF files (1-based page numbers).

    Uses PyMuPDF to extract text page by page.  Tables are detected
    automatically and rendered as Markdown to preserve structure.
    """
    try:
        doc = fitz.open(path)
    except Exception as exc:
        LOGGER.error("Failed to open PDF %s: %s", path, exc)
        return
    try:
        for index in range(len(doc)):
            try:
                page = doc[index]
                text = _extract_page_text(page)
                normalized = normalize_whitespace([text])
                if normalized:
                    yield index + 1, normalized + "\n"
            except Exception as exc:  # pragma: no cover
                LOGGER.warning("Failed to read page %s in %s: %s", index, path, exc)
    finally:
        doc.close()


def _iter_paged_text(path: Path) -> Iterator[tuple[int, str]]:
    """Yield ``(page_number, text)`` for any supported format.

    * PDF  → real page numbers (1-based)
    * Markdown → section numbers (split on headings)
    * DOCX / ODF → virtual pages (every 10 paragraphs)
    * PPTX → slide numbers
    * EPUB → chapter numbers
    * Plain text / HTML / DOC / PPT → virtual pages (~3000 characters)
    """
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        yield from iter_text_parts_paged(path)
    elif suffix == ".md":
        yield from iter_text_parts_md_paged(path)
    elif suffix == ".docx":
        yield from iter_text_parts_docx_paged(path)
    elif suffix == ".txt":
        yield from iter_text_parts_txt_paged(path)
    elif suffix in (".odt", ".odp", ".odg"):
        yield from iter_text_parts_odf_paged(path)
    elif suffix == ".pptx":
        yield from iter_text_parts_pptx_paged(path)
    elif suffix in (".html", ".htm"):
        yield from iter_text_parts_html_paged(path)
    elif suffix == ".epub":
        yield from iter_text_parts_epub_paged(path)
    elif suffix == ".doc":
        yield from iter_text_parts_doc_paged(path)
    elif suffix == ".ppt":
        yield from iter_text_parts_ppt_paged(path)
    else:
        LOGGER.warning("Unsupported file type: %s", path.suffix)


def build_chunks(path: Path, *, max_chars: int = 1200, overlap: int = 200) -> Iterable[ChunkRecord]:
    """Produce overlapping chunk records for any supported document type."""
    from docfinder.utils.text import chunk_text_stream_paged

    title = _get_title(path)
    pages = _iter_paged_text(path)

    for idx, (chunk, page_num) in enumerate(
        chunk_text_stream_paged(pages, max_chars=max_chars, overlap=overlap)
    ):
        yield ChunkRecord(
            document_path=path,
            index=idx,
            text=chunk,
            metadata={"title": title, "page": page_num},
        )
